from pptx.util import Inches

from pptx import Presentation

# Load the existing presentation
prs = Presentation("Efficient_Application_Classification.pptx")

# Add a new slide with a title and content layout
slide_layout = prs.slide_layouts[1]  # Title and content layout
slide = prs.slides.add_slide(slide_layout)

# Set the slide title
title = slide.shapes.title
title.text = "Model Training and Evaluation"

# Set the content
content = (
    "Random Forest on Full PCA-Reduced Data:\n"
    "- Trained a Random Forest model on the PCA-reduced data.\n"
    "- Achieved high accuracy on the test set.\n\n"
    "Random Forest on Core Set Data:\n"
    "- Trained a model on the selected core set.\n"
    "- Core set provides a compressed version of the full data.\n"
    "- Evaluated model performance on both core set and full PCA-reduced data."
)

# Add the content text box
text_box = slide.shapes.placeholders[1].text = content

# Add a code snippet to the slide
left = Inches(0.5)
top = Inches(4.0)
width = Inches(8.5)
height = Inches(2.5)
code_snippet = (
    "# Training Random Forest on Full Data\n"
    "rf_full = RandomForestClassifier(n_estimators=50, n_jobs=-1, random_state=42)\n"
    "rf_full.fit(X_train_full, y_train_full)\n"
    "accuracy_full = accuracy_score(y_test_full, y_pred_full)\n\n"
    "# Training Random Forest on Core Set Data\n"
    "rf_core = RandomForestClassifier(n_estimators=100, n_jobs=-1, random_state=42)\n"
    "rf_core.fit(X_train_core, y_train_core)\n"
    "accuracy_core = accuracy_score(y_test_core, y_pred_core)\n\n"
    "# Evaluation on Full Data using Core Set Model\n"
    "accuracy_full_core = accuracy_score(y, y_pred_full_core)"
)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.text = code_snippet

# Save the updated presentation
prs.save("Efficient_Application_Classification.pptx")
